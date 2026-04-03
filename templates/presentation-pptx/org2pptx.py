#!/usr/bin/env python3
"""org2pptx — Inject org-mode content into a PowerPoint template.

Parses an org file where each level-1 heading maps to a slide in the
template, then injects bullets, tables and sub-sections via python-pptx
while preserving the original template design (backgrounds, logos, layouts).

Architecture:
    [content.org] → [org parser] → [slide dicts] → [python-pptx] → [output.pptx]
                                                         ↑
                                                  [template.pptx]

Org conventions:
    - Level 1 heading  = one slide
    - :SLIDE_INDEX: N  = target slide number in template (1-based)
    - :LAYOUT:         = informational only (template already has it)
    - :SKIP: true      = keep template slide as-is (no injection)
    - :REPLACE_TEXT: old=new  = replace text in existing shapes
    - Level 2 heading  = sub-section (▶ title + bullets / table)
    - Bullet list      = ``- item`` lines
    - Org table        = ``| col | col |`` lines (separator ``|---`` skipped)

Requires: orgparse, python-pptx
    nix-shell -p python312Packages.orgparse python312Packages.python-pptx

Examples:
    python org2pptx.py slides.org -t template.pptx
    python org2pptx.py slides.org -t template.pptx -o output.pptx
    python org2pptx.py slides.org -t template.pptx --remove-decorations "직사각형"
"""
import argparse
import shutil
import sys
from dataclasses import dataclass, field
from pathlib import Path

import orgparse
from pptx import Presentation
from pptx.util import Cm, Pt


# ── Org parsing ───────────────────────────────────────────

@dataclass
class SubSection:
    """Level 2 heading: sub-section within a slide."""
    title: str
    bullets: list[str] = field(default_factory=list)
    table: list[list[str]] = field(default_factory=list)


@dataclass
class SlideContent:
    """Level 1 heading: one slide."""
    title: str
    slide_index: int  # 1-based
    layout: str = ""
    skip: bool = False
    replace_text: dict[str, str] = field(default_factory=dict)
    bullets: list[str] = field(default_factory=list)
    table: list[list[str]] = field(default_factory=list)
    subsections: list[SubSection] = field(default_factory=list)


def parse_body(body: str) -> tuple[list[str], list[list[str]]]:
    """Extract bullet list and table from org body text."""
    bullets = []
    table = []
    for line in body.split("\n"):
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        if line.startswith("- "):
            bullets.append(line[2:].strip())
        elif line.startswith("| ") and not line.startswith("|---"):
            cells = [c.strip() for c in line.split("|")[1:-1]]
            table.append(cells)
    return bullets, table


def parse_replace_text(raw: str) -> dict[str, str]:
    """Parse :REPLACE_TEXT: 'old=new' into a dict.

    Supports multiple replacements separated by ``;;``.
    Example: ``발표기관명 작성=ACME Corp;;2024=2025``
    """
    replacements = {}
    if not raw:
        return replacements
    for pair in raw.split(";;"):
        pair = pair.strip()
        if "=" in pair:
            old, new = pair.split("=", 1)
            replacements[old.strip()] = new.strip()
    return replacements


def parse_org(org_path: str) -> list[SlideContent]:
    """Parse org file into a list of SlideContent."""
    root = orgparse.load(org_path)
    slides = []

    for node in root.children:
        if node.level != 1:
            continue

        props = node.properties
        slide_index = int(props.get("SLIDE_INDEX", 0))
        if slide_index == 0:
            print(f"  ⚠ [{node.heading}]: no SLIDE_INDEX, skipping",
                  file=sys.stderr)
            continue

        skip = props.get("SKIP", "").lower() == "true"
        layout = props.get("LAYOUT", "")
        replace_text = parse_replace_text(props.get("REPLACE_TEXT", ""))

        bullets, table = parse_body(node.body)

        subsections = []
        for child in node.children:
            if child.level == 2:
                sub_bullets, sub_table = parse_body(child.body)
                subsections.append(SubSection(
                    title=child.heading,
                    bullets=sub_bullets,
                    table=sub_table,
                ))

        slides.append(SlideContent(
            title=node.heading,
            slide_index=slide_index,
            layout=layout,
            skip=skip,
            replace_text=replace_text,
            bullets=bullets,
            table=table,
            subsections=subsections,
        ))

    return slides


# ── PPTX injection ────────────────────────────────────────

def get_slides(prs: Presentation) -> list:
    """Return slide list (workaround for broken index access in some pptx)."""
    slides = []
    for sldId in prs.slides._sldIdLst:
        slides.append(prs.slides.part.related_slide(sldId.rId))
    return slides


def add_textbox(slide, left_cm, top_cm, width_cm, height_cm):
    """Add a textbox and return its TextFrame."""
    txBox = slide.shapes.add_textbox(
        Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def write_subsection_title(tf, title: str, is_first: bool = False):
    """Write a sub-section title (▶ prefix, bold)."""
    if is_first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
        p.space_before = Pt(12)
    p.text = f"▶ {title}"
    for run in p.runs:
        run.font.size = Pt(16)
        run.font.bold = True


def write_bullets(tf, bullets: list[str]):
    """Write bullet items to a TextFrame."""
    for item in bullets:
        p = tf.add_paragraph()
        p.text = f"  • {item}"
        for run in p.runs:
            run.font.size = Pt(13)
        p.space_before = Pt(3)


def write_table(slide, table: list[list[str]], top_cm: float):
    """Write an org table as a pptx Table shape."""
    if not table:
        return
    rows = len(table)
    cols = len(table[0])
    tbl_shape = slide.shapes.add_table(
        rows, cols,
        Cm(2.5), Cm(top_cm), Cm(22), Cm(min(rows * 1.5, 12))
    )
    tbl = tbl_shape.table
    for i, row in enumerate(table):
        for j, cell_text in enumerate(row):
            if j < len(tbl.columns):
                tbl.cell(i, j).text = cell_text


def remove_decoration_shapes(slide, prefixes: list[str]) -> int:
    """Remove non-placeholder shapes whose name starts with any prefix.

    Some templates have decorative shapes (e.g. colored rectangles with
    duplicated slide titles) that overlap with injected content.
    This function removes them by name prefix.

    Args:
        slide: pptx slide object
        prefixes: list of shape name prefixes to match (e.g. ["직사각형"])

    Returns:
        Number of shapes removed.
    """
    if not prefixes:
        return 0
    sp_tree = slide.shapes._spTree
    to_remove = []
    for shape in slide.shapes:
        if (not shape.is_placeholder
                and shape.has_text_frame
                and any(shape.name.startswith(p) for p in prefixes)):
            to_remove.append(shape._element)
    for elem in to_remove:
        sp_tree.remove(elem)
    return len(to_remove)


def replace_shape_text(slide, replacements: dict[str, str]) -> int:
    """Replace text in existing shapes (run-level for style preservation).

    Args:
        slide: pptx slide object
        replacements: {old_text: new_text} dict

    Returns:
        Number of replacements made.
    """
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        count += 1
    return count


def inject_slide(pptx_slide, content: SlideContent,
                 deco_prefixes: list[str] | None = None):
    """Inject SlideContent into a pptx slide."""
    has_subsections = bool(content.subsections)
    has_direct_content = bool(content.bullets or content.table)

    if not has_subsections and not has_direct_content:
        return

    # Remove decoration shapes that would overlap with content
    if deco_prefixes:
        remove_decoration_shapes(pptx_slide, deco_prefixes)

    text_top = 4.0  # cm below slide title area

    if has_subsections:
        text_subs = [s for s in content.subsections if s.bullets]
        table_subs = [s for s in content.subsections if s.table]

        if text_subs:
            tf = add_textbox(pptx_slide, 2.5, text_top, 22, 13)
            first = True
            for sub in text_subs:
                write_subsection_title(tf, sub.title, is_first=first)
                first = False
                write_bullets(tf, sub.bullets)

        tbl_top = text_top + 1.5 + sum(
            len(s.bullets) * 0.7 + 1.2 for s in text_subs
        )
        for sub in table_subs:
            write_table(pptx_slide, sub.table, min(tbl_top, 12.0))
            tbl_top += len(sub.table) * 1.2 + 1.0

    elif has_direct_content:
        if content.bullets:
            tf = add_textbox(pptx_slide, 2.5, text_top, 22, 13)
            p = tf.paragraphs[0]
            p.text = f"  • {content.bullets[0]}"
            for run in p.runs:
                run.font.size = Pt(13)
            write_bullets(tf, content.bullets[1:])

        if content.table:
            table_offset = text_top + (
                len(content.bullets) * 0.7 if content.bullets else 0
            )
            write_table(pptx_slide, content.table, min(table_offset, 5.0))


# ── Main pipeline ─────────────────────────────────────────

def org2pptx(org_path: str, template_path: str, output_path: str,
             deco_prefixes: list[str] | None = None):
    """Main pipeline: org → parse → inject → save."""
    # 1. Parse org
    print(f"📖 Parsing: {org_path}")
    slides_content = parse_org(org_path)
    print(f"   {len(slides_content)} slide(s) defined")

    # 2. Copy template
    shutil.copy(template_path, output_path)
    prs = Presentation(output_path)
    pptx_slides = get_slides(prs)
    print(f"📄 Template: {len(pptx_slides)} slide(s) ({template_path})")

    # 3. Inject
    for sc in slides_content:
        idx = sc.slide_index - 1
        if idx < 0 or idx >= len(pptx_slides):
            print(f"  ⚠ [{sc.title}]: SLIDE_INDEX {sc.slide_index} "
                  f"out of range (1-{len(pptx_slides)})", file=sys.stderr)
            continue

        pptx_slide = pptx_slides[idx]

        # Text replacement (works on both SKIP and non-SKIP slides)
        if sc.replace_text:
            n = replace_shape_text(pptx_slide, sc.replace_text)
            print(f"  🔄 Slide {sc.slide_index}: {n} text replacement(s)")

        if sc.skip:
            if not sc.replace_text:
                print(f"  ⏭ Slide {sc.slide_index}: SKIP ({sc.title})")
            continue

        inject_slide(pptx_slide, sc, deco_prefixes)
        n_sub = len(sc.subsections)
        n_bul = len(sc.bullets) + sum(
            len(s.bullets) for s in sc.subsections
        )
        n_tbl = (1 if sc.table else 0) + sum(
            1 for s in sc.subsections if s.table
        )
        deco = ", deco_removed" if deco_prefixes else ""
        print(f"  ✅ Slide {sc.slide_index}: {sc.title} "
              f"(sub={n_sub}, bullets={n_bul}, tables={n_tbl}{deco})")

    # 4. Save
    prs.save(output_path)
    print(f"\n💾 Saved: {output_path}")


# ── CLI ───────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Inject org-mode content into a PowerPoint template",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""\
examples:
  %(prog)s slides.org -t template.pptx
  %(prog)s slides.org -t template.pptx -o final.pptx
  %(prog)s slides.org -t template.pptx --remove-decorations "직사각형"
""",
    )
    parser.add_argument("org_file", help="Input org file")
    parser.add_argument(
        "--template", "-t", required=True,
        help="PowerPoint template file (.pptx)"
    )
    parser.add_argument(
        "--output", "-o", default=None,
        help="Output pptx path (default: <org_stem>.pptx)"
    )
    parser.add_argument(
        "--remove-decorations", "-r", default=None,
        help="Comma-separated shape name prefixes to remove "
             "(e.g. '직사각형' for Korean templates with decoration rects)"
    )
    args = parser.parse_args()

    if args.output is None:
        args.output = Path(args.org_file).stem + ".pptx"

    deco_prefixes = None
    if args.remove_decorations:
        deco_prefixes = [p.strip() for p in args.remove_decorations.split(",")]

    org2pptx(args.org_file, args.template, args.output, deco_prefixes)


if __name__ == "__main__":
    main()
